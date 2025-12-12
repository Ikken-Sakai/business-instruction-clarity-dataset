#!/usr/bin/env python3
"""
中間発表 PowerPoint資料 生成スクリプト
外国人労働者のための日本語ビジネス指示文の曖昧性判定システム
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

def set_text_frame_properties(text_frame, font_size=18, font_bold=False, font_color=None, alignment=PP_ALIGN.LEFT):
    """テキストフレームのプロパティを設定"""
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = alignment
        for run in paragraph.runs:
            run.font.size = Pt(font_size)
            run.font.bold = font_bold
            if font_color:
                run.font.color.rgb = font_color

def add_title_slide(prs):
    """スライド1: タイトル"""
    slide_layout = prs.slide_layouts[6]  # 空白レイアウト
    slide = prs.slides.add_slide(slide_layout)
    
    # メインタイトル
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "外国人労働者のための\n日本語ビジネス指示文の曖昧性判定システム"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    
    # サブタイトル
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.8))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "〜 BERT二値分類モデルによる明確性自動判定 〜"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(80, 80, 80)
    
    # 発表者情報
    info_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1))
    tf = info_box.text_frame
    p = tf.paragraphs[0]
    p.text = "発表開始時刻　○○:○○\n○○_TK230422_坂井壱謙（さかいいっけん）"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(60, 60, 60)

def add_background_slide(prs):
    """スライド2: 背景・目的"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # スライドタイトル
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "背景・目的"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    
    # 背景セクション
    bg_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(0.4))
    tf = bg_title.text_frame
    p = tf.paragraphs[0]
    p.text = "【背景】"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 153)
    
    bg_content = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(2))
    tf = bg_content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "■ 日本の職場における曖昧な指示"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "　　「例の件、早めに対応しといて」"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(100, 100, 100)
    
    p = tf.add_paragraph()
    p.text = "　　「いい感じにまとめといて」"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(100, 100, 100)
    
    p = tf.add_paragraph()
    p.text = ""
    p.font.size = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "■ 外国人労働者が感じる困難"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "　　→「何を」「いつまでに」が不明確"
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "　　→ ミスコミュニケーションの原因（安部, 2018）"
    p.font.size = Pt(14)
    
    # 目的セクション
    purpose_title = slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(9), Inches(0.4))
    tf = purpose_title.text_frame
    p = tf.paragraphs[0]
    p.text = "【目的】"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 153)
    
    purpose_content = slide.shapes.add_textbox(Inches(0.5), Inches(4.4), Inches(9), Inches(1.2))
    tf = purpose_content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "BERTを用いて指示文の「明確/曖昧」を自動判定"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "→ 曖昧な指示をリアルタイムで検出・改善支援"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(0, 102, 153)

def add_method_slide(prs):
    """スライド3: 方法"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # スライドタイトル
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "方法"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    
    # 手法セクション
    method_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(0.4))
    tf = method_title.text_frame
    p = tf.paragraphs[0]
    p.text = "【手法】"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 153)
    
    method_content = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1.5))
    tf = method_content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "モデル: BERT (cl-tohoku/bert-base-japanese-v3)"
    p.font.size = Pt(16)
    
    p = tf.add_paragraph()
    p.text = "タスク: 二値分類（明確 vs 曖昧）"
    p.font.size = Pt(16)
    
    p = tf.add_paragraph()
    p.text = ""
    p.font.size = Pt(8)
    
    p = tf.add_paragraph()
    p.text = "入力: ビジネス指示文  →  出力: 明確 / 曖昧"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 153)
    
    # データセットセクション
    data_title = slide.shapes.add_textbox(Inches(0.5), Inches(3.3), Inches(9), Inches(0.4))
    tf = data_title.text_frame
    p = tf.paragraphs[0]
    p.text = "【学習データセット】"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 153)
    
    data_content = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(9), Inches(1.8))
    tf = data_content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "・総データ数: 2,000件（明確:曖昧 = 50:50）"
    p.font.size = Pt(16)
    
    p = tf.add_paragraph()
    p.text = "・分割: Train 1,600 / Val 200 / Test 200"
    p.font.size = Pt(16)
    
    p = tf.add_paragraph()
    p.text = "・10業種 × 多様なシチュエーション × 3種類の口調"
    p.font.size = Pt(16)

def add_result_slide(prs):
    """スライド4: 結果（進捗報告）"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # スライドタイトル
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "結果（進捗報告）"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    
    # 進捗状況セクション
    progress_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(0.4))
    tf = progress_title.text_frame
    p = tf.paragraphs[0]
    p.text = "【進捗状況】"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 153)
    
    # 完了項目
    done_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(4), Inches(2))
    tf = done_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "✅ 完了"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 128, 0)
    
    p = tf.add_paragraph()
    p.text = "・データセット構築（2,000件）"
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "・アノテーション基準策定"
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "・品質チェック完了"
    p.font.size = Pt(14)
    
    # 未着手項目
    todo_box = slide.shapes.add_textbox(Inches(5), Inches(1.6), Inches(4.5), Inches(2))
    tf = todo_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "⬜ 未着手"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(128, 128, 128)
    
    p = tf.add_paragraph()
    p.text = "・BERTモデル学習"
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "・モデル評価"
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "・デモアプリ開発"
    p.font.size = Pt(14)
    
    # データ統計セクション
    stats_title = slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(9), Inches(0.4))
    tf = stats_title.text_frame
    p = tf.paragraphs[0]
    p.text = "【データ統計】"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 153)
    
    stats_content = slide.shapes.add_textbox(Inches(0.5), Inches(4.4), Inches(9), Inches(1))
    tf = stats_content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "・文字数: 10〜49字（平均24.2字）"
    p.font.size = Pt(16)
    
    p = tf.add_paragraph()
    p.text = "・ラベルバランス: 完全に50:50"
    p.font.size = Pt(16)

def add_summary_slide(prs):
    """スライド5: まとめ・考察"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # スライドタイトル
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "まとめ・考察"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    
    # 現状まとめ
    summary_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(0.4))
    tf = summary_title.text_frame
    p = tf.paragraphs[0]
    p.text = "【現状まとめ】"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 153)
    
    summary_content = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(0.8))
    tf = summary_content.text_frame
    p = tf.paragraphs[0]
    p.text = "・高品質なデータセット構築完了"
    p.font.size = Pt(14)
    p = tf.add_paragraph()
    p.text = "・モデル学習は次フェーズで実施"
    p.font.size = Pt(14)
    
    # 予想結果
    expect_title = slide.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(9), Inches(0.4))
    tf = expect_title.text_frame
    p = tf.paragraphs[0]
    p.text = "【予想される結果】"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 153)
    
    expect_content = slide.shapes.add_textbox(Inches(0.5), Inches(2.7), Inches(9), Inches(0.5))
    tf = expect_content.text_frame
    p = tf.paragraphs[0]
    p.text = "・目標: F1スコア 85%以上"
    p.font.size = Pt(14)
    
    # 本番までの計画
    plan_title = slide.shapes.add_textbox(Inches(0.5), Inches(3.3), Inches(9), Inches(0.4))
    tf = plan_title.text_frame
    p = tf.paragraphs[0]
    p.text = "【本番までの計画】"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 153)
    
    # 計画テーブル風
    plan_content = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(9), Inches(1.8))
    tf = plan_content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "12月中旬　　　│　　12月下旬　　　│　　1月上旬"
    p.font.size = Pt(14)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "─────────────┼─────────────────┼───────────────"
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "BERT学習　　　│　　評価・分析　　│　　デモアプリ"
    p.font.size = Pt(14)
    
    p = tf.add_paragraph()
    p.text = "パラメータ調整│　　結果考察　　　│　　発表準備"
    p.font.size = Pt(14)

def add_appendix_label_slide(prs):
    """スライド6: APPENDIX - ラベル定義"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # スライドタイトル
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "APPENDIX: ラベル定義詳細"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    
    # Label 0
    label0_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(0.4))
    tf = label0_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Label 0（明確）: What + When が両方具体的"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 128, 0)
    
    label0_example = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.5), Inches(0.8))
    tf = label0_example.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "例:「今日17時までにA社の見積書をPDFで作成して」"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(80, 80, 80)
    
    p = tf.add_paragraph()
    p.text = "　　→ 期限・対象・形式が全て明確"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(100, 100, 100)
    
    # Label 1
    label1_title = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(0.4))
    tf = label1_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Label 1（曖昧）: What または When が欠如"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(200, 0, 0)
    
    label1_example = slide.shapes.add_textbox(Inches(0.7), Inches(2.9), Inches(8.5), Inches(0.8))
    tf = label1_example.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "例:「例の件、早めに対応しといて」"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(80, 80, 80)
    
    p = tf.add_paragraph()
    p.text = "　　→「例の件」が不明、「早めに」が主観的"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(100, 100, 100)
    
    # 判定基準表
    table_title = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.4))
    tf = table_title.text_frame
    p = tf.paragraphs[0]
    p.text = "【判定基準】"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 153)
    
    table_content = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1.5))
    tf = table_content.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "項目　　　│　曖昧（×）　　　　　│　明確（○）"
    p.font.size = Pt(12)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "────────┼────────────────────┼────────────────────"
    p.font.size = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "動詞　　　│　対応する、処理する　│　メール送る、印刷する"
    p.font.size = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "対象　　　│　例の件、あの資料　　│　A社の請求書"
    p.font.size = Pt(12)
    
    p = tf.add_paragraph()
    p.text = "期限　　　│　なる早、早めに　　　│　15時までに、今日中に"
    p.font.size = Pt(12)

def add_appendix_reference_slide(prs):
    """スライド7: APPENDIX - 参考文献"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # スライドタイトル
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "APPENDIX: 参考文献"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    
    # 参考文献リスト
    ref_content = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(4.5))
    tf = ref_content.text_frame
    tf.word_wrap = True
    
    references = [
        "1. 安部陽子 (2018). 「日本での就労時に留学生が持つ違和感の調査報告：\n　　日本人学生との対照分析を通して」, 地球社会統合科学研究, 第8号, pp.1-12.",
        "",
        "2. 宗甜甜 (2021). 「ビジネス場面における日本語の『断り』に関する研究」,\n　　博士論文（日本大学）.",
        "",
        "3. 東中竜一郎 他 (2016). 「テキストチャットを用いた雑談対話コーパスの\n　　構築と対話破綻の分析」, 自然言語処理, Vol.23, No.1.",
        "",
        "4. 奥村学, 田中穂積 (1989). 「自然言語解析における意味的曖昧性を増進的に\n　　解消する計算モデル」, 人工知能学会誌, Vol.4, No.6.",
        "",
        "5. Devlin, J., Chang, M. W., Lee, K., & Toutanova, K. (2018).\n　　\"BERT: Pre-training of Deep Bidirectional Transformers for\n　　Language Understanding.\" arXiv preprint arXiv:1810.04805.",
    ]
    
    p = tf.paragraphs[0]
    p.text = references[0]
    p.font.size = Pt(12)
    
    for ref in references[1:]:
        p = tf.add_paragraph()
        p.text = ref
        p.font.size = Pt(12)

def main():
    """メイン処理"""
    # プレゼンテーション作成
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # 各スライドを追加
    add_title_slide(prs)
    add_background_slide(prs)
    add_method_slide(prs)
    add_result_slide(prs)
    add_summary_slide(prs)
    add_appendix_label_slide(prs)
    add_appendix_reference_slide(prs)
    
    # 保存
    output_path = os.path.join(os.path.dirname(__file__), "中間発表_外国人労働者向け曖昧性判定システム.pptx")
    prs.save(output_path)
    print(f"✅ PowerPoint発表資料を生成しました: {output_path}")

if __name__ == "__main__":
    main()

